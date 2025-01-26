import os
import pdfplumber
from pptx import Presentation

def save_uploaded_file(file):
    """
    Save a single uploaded file and extract text based on its type.
    """
    if file.filename == '':
        raise ValueError("No file provided")

    # Save the file in the upload folder
    upload_folder = os.path.join('app/uploads')
    os.makedirs(upload_folder, exist_ok=True)

    file_path = os.path.join(upload_folder, file.filename)
    file.save(file_path)

    # Extract text based on file type
    if file.filename.lower().endswith('.pdf'):
        content = extract_text_from_pdf(file_path)
    elif file.filename.lower().endswith('.txt'):
        content = extract_text_from_txt(file_path)
    elif file.filename.lower().endswith('.pptx'):
        content = extract_text_from_ppt(file_path)
    else:
        raise ValueError("Unsupported file type. Only PDF, TXT, and PPTX files are supported.")

    return file.filename, content

def extract_text_from_pdf(file_path):
    """
    Extract text from a PDF file using pdfplumber.
    """
    try:
        text = ''
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() + '\n'
        return text.strip()
    except Exception as e:
        print(f"Error reading PDF {file_path}: {e}")
        return 'Error extracting text from this file.'

def extract_text_from_txt(file_path):
    """
    Extract text from a TXT file.
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as txt_file:
            return txt_file.read().strip()
    except Exception as e:
        print(f"Error reading TXT file {file_path}: {e}")
        return 'Error extracting text from this file.'

def extract_text_from_ppt(file_path):
    """
    Extract text from a PPTX file using python-pptx.
    """
    try:
        presentation = Presentation(file_path)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        text.append(paragraph.text)
        return "\n".join(text).strip()
    except Exception as e:
        print(f"Error reading PPTX file {file_path}: {e}")
        return 'Error extracting text from this file.'
