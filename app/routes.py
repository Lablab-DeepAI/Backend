from flask import Blueprint, request, jsonify, current_app
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import pandas as pd
import os
import pdfplumber
from pptx import Presentation  # To handle PowerPoint files
from app.chatbot import ask_groq

main = Blueprint('main', __name__)

# Dictionary to store uploaded content in memory
uploaded_file_content = {}

vectorizer = TfidfVectorizer(stop_words='english')

@main.route('/upload', methods=['POST'])
def upload_file():
    """
    Endpoint to upload a file (PDF, PPT, or TXT) and extract its content.
    """
    if 'file' not in request.files:
        return jsonify({'error': 'No file part in the request'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400

    # Check file extension
    allowed_extensions = ['.pdf', '.pptx', '.txt']
    file_ext = os.path.splitext(file.filename)[1].lower()

    if file and file_ext in allowed_extensions:
        upload_folder = current_app.config['UPLOAD_FOLDER']
        os.makedirs(upload_folder, exist_ok=True)
        filepath = os.path.join(upload_folder, file.filename)
        file.save(filepath)

        try:
            if file_ext == '.pdf':
                # Extract text from PDF
                with pdfplumber.open(filepath) as pdf:
                    text = "\n".join([page.extract_text() for page in pdf.pages if page.extract_text()])
            elif file_ext == '.pptx':
                # Extract text from PPT
                ppt = Presentation(filepath)
                text = "\n".join([paragraph.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text_frame") for paragraph in shape.text_frame.paragraphs])
            elif file_ext == '.txt':
                # Extract text from TXT
                with open(filepath, 'r', encoding='utf-8') as txt_file:
                    text = txt_file.read()

            # Store the content in memory for future use
            uploaded_file_content[file.filename] = text

            return jsonify({'message': 'File uploaded successfully', 'filename': file.filename, 'content':text}), 200

        except Exception as e:
            return jsonify({'error': f"Failed to process the file: {str(e)}"}), 500

    return jsonify({'error': f'Invalid file type, only {allowed_extensions} are allowed'}), 400


@main.route('/chat', methods=['POST'])
def chat():
    """
    Endpoint to handle chatbot queries.
    Expects JSON input with 'question' and 'filename' keys.
    """
    data = request.json
    question = data.get('question')
    filename = data.get('filename')

    if not question or not filename:
        return jsonify({"error": "Both 'question' and 'filename' are required"}), 400

    # Retrieve the content of the uploaded file from the file itself
    file_path = os.path.join(current_app.config['UPLOAD_FOLDER'], filename)
    if not os.path.exists(file_path):
        return jsonify({"error": f"File '{filename}' not found"}), 404

    try:
        # Check if content is already in memory
        if filename in uploaded_file_content:
            content = uploaded_file_content[filename]
        else:
            # Reload content from the file if not in memory
            file_ext = os.path.splitext(filename)[1].lower()
            if file_ext == '.pdf':
                with pdfplumber.open(file_path) as pdf:
                    content = "\n".join([page.extract_text() for page in pdf.pages if page.extract_text()])
            elif file_ext == '.pptx':
                ppt = Presentation(file_path)
                content = "\n".join([paragraph.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text_frame") for paragraph in shape.text_frame.paragraphs])
            elif file_ext == '.txt':
                with open(file_path, 'r', encoding='utf-8') as txt_file:
                    content = txt_file.read()
            else:
                return jsonify({'error': 'Unsupported file type'}), 400

        # Use the chatbot function to get an answer
        groq_response = ask_groq(question, content)

        if "error" in groq_response:
            return jsonify({"error": groq_response["error"]}), 500

        return jsonify({
            "answer": groq_response.get("answer"),
            "confidence": groq_response.get("confidence")
        })

    except Exception as e:
        return jsonify({'error': f"Failed to process the file: {str(e)}"}), 500

#Endpoint accepts a user id and returns recommended documents based on the user's activity
@main.route('/recommend_from_user', methods=['POST'])
def recommend_from_user(): 
    user_id = request.json['user_id']
    #To Do: get user activity from db for documents the user recently interacted with in db
    #compare these documents to documents in db and return ones with highest similarity 

#Endpoint accepts a document and recommends documents similar to the provided document
@main.route('/recommend_from_document', methods=['POST'])
def recommend_from_document():
    if 'document_id' in request.json:
        try: 
            document_id = request.json['document_id']
            #To Do: get content for all documents in the db and store in pandas dataframe 
            #To Do: specify query string in read_sql function below once db structure is defined
            documents_df = pd.read_sql()
            tfidf_matrix = vectorizer.fit_transform(documents_df['content'])
            # To Do: get content from db for the document with given id 
            # To Do: replace str below with the result from db
            document_content = str 
            input_document_vector = vectorizer.transform([document_content])
            cosine_similarities = cosine_similarity(input_document_vector, tfidf_matrix)
            cosine_sim_df = pd.DataFrame(cosine_similarities.T, columns=["similarity"], index=documents_df["title"])
            similar_documents = cosine_sim_df.sort_values(by="similarity", ascending=False)
            n = 10
            top_similar_movies = similar_documents.head(n)
            return top_similar_movies
        except:
            return jsonify({'error': 'Could not find recommendations'}), 400
    return jsonify({'error': 'Document Id was not provided'}), 400